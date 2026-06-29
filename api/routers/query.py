import base64
from fastapi import APIRouter, Depends, HTTPException, BackgroundTasks, Request, UploadFile, File, Form, Response
from pydantic import BaseModel, Field
from typing import Optional
import structlog

from api.auth_utils import get_current_user_or_api_key
from src.agent import ask
from src.memory.sqlite_store import write_audit
from src.generation.llm import MODEL_REGISTRY, MODEL_METADATA

logger = structlog.get_logger()
router = APIRouter()


# ── Free-trial limit enforcement ──────────────────────────────────────────────

def _user_email(user_id: str) -> str:
    """Look up a user's email (JWT only carries user_id, not email)."""
    from src.memory.sqlite_store import get_conn
    try:
        with get_conn() as c:
            row = c.execute("SELECT email FROM users WHERE user_id=?", (user_id,)).fetchone()
        return (row[0] or "").lower() if row else ""
    except Exception:
        return ""


def _user_question_count(user_id: str) -> int:
    """Lifetime count of metered questions (audit_log 'query' events) for a user."""
    from src.memory.sqlite_store import get_conn
    try:
        with get_conn() as c:
            row = c.execute(
                "SELECT COUNT(*) FROM audit_log WHERE event_type='query' AND user_id=?",
                (user_id,),
            ).fetchone()
        return int(row[0]) if row else 0
    except Exception:
        return 0


def _enforce_free_limit(user: dict) -> None:
    """Block non-owner free users once they hit the lifetime question limit.

    Owners (OWNER_EMAILS) and paid plans bypass. Raises HTTP 402 when exceeded.
    """
    from src.config import OWNER_EMAILS, FREE_QUESTION_LIMIT
    from api.routers.billing import get_user_plan
    email = _user_email(user["user_id"])
    if email in OWNER_EMAILS:
        return  # owner / co-owner → unlimited
    plan, _ = get_user_plan(user["user_id"])
    if plan != "free":
        return  # active paid plan → unlimited
    used = _user_question_count(user["user_id"])
    if used >= FREE_QUESTION_LIMIT:
        raise HTTPException(
            status_code=402,
            detail=(f"Free trial limit reached — you've used all {FREE_QUESTION_LIMIT} "
                    f"free questions. Please upgrade to continue."),
        )


@router.get("/usage")
async def usage(user: dict = Depends(get_current_user_or_api_key)):
    """Return the current user's plan and free-trial usage for the Plan page."""
    from src.config import OWNER_EMAILS, FREE_QUESTION_LIMIT
    from api.routers.billing import get_user_plan
    email = _user_email(user["user_id"])
    is_owner = email in OWNER_EMAILS
    plan, expires_at = get_user_plan(user["user_id"])
    unlimited = is_owner or plan != "free"
    used = _user_question_count(user["user_id"])
    return {
        "plan": "owner" if is_owner else plan,
        "unlimited": unlimited,
        "used": used,
        "limit": None if unlimited else FREE_QUESTION_LIMIT,
        "remaining": None if unlimited else max(0, FREE_QUESTION_LIMIT - used),
        "expires_at": expires_at,
    }


class QueryRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=4000)
    session_id: Optional[str] = Field(default="default", max_length=64)
    stream: bool = False
    language: str = Field(default="English", max_length=50)
    source_filters: list[str] = Field(default_factory=list)
    force_bi: bool = False
    model: Optional[str] = Field(default=None, max_length=100)  # api_model_id from MODEL_REGISTRY
    no_cache: bool = False  # force a fresh answer, bypassing the verified-knowledge cache


class TokenUsage(BaseModel):
    input_tokens: int = 0
    output_tokens: int = 0
    total_tokens: int = 0
    llm_calls: int = 0
    estimated_cost_usd: float = 0.0
    model: str = ""  # model that actually produced the answer (reflects fallbacks)


class QueryResponse(BaseModel):
    answer: str
    quality_score: float
    confidence: str
    faithfulness: str
    patterns_used: list[str]
    latency_ms: int
    retrieval_channel: str
    fallback_used: bool
    verified_knowledge_hit: bool
    suggested_followups: list[str]
    citation_map: dict
    figures: list[str] = []
    session_id: str
    model_used: Optional[str] = None
    token_usage: Optional[TokenUsage] = None
    request_id: Optional[str] = None


@router.get("/models")
async def list_models(user: dict = Depends(get_current_user_or_api_key)):
    """Return available LLM models with pricing/context metadata for the UI."""
    models = []
    for label, (provider, model_id) in MODEL_REGISTRY.items():
        meta = MODEL_METADATA.get(model_id, {})
        models.append({
            "label": label,
            "model_id": model_id,
            "provider": provider,
            "context_k": meta.get("context_k", 0),
            "input_usd_per_mtok": meta.get("input_usd_per_mtok", 0.0),
            "output_usd_per_mtok": meta.get("output_usd_per_mtok", 0.0),
            "free": meta.get("free", False),
        })
    return {"models": models}


@router.get("/token-stats")
async def token_stats(user: dict = Depends(get_current_user_or_api_key)):
    """Return cumulative token usage for the current user from audit logs."""
    from src.memory.sqlite_store import get_conn
    conn = get_conn()
    try:
        rows = conn.execute(
            "SELECT detail FROM audit_log WHERE event_type='query' AND user_id=?",
            (user["user_id"],),
        ).fetchall()
    finally:
        conn.close()

    total_input = 0
    total_output = 0
    total_cost = 0.0
    query_count = 0

    for (detail_str,) in rows:
        try:
            detail = json.loads(detail_str) if detail_str else {}
        except Exception:
            continue
        tu = detail.get("token_usage") or {}
        total_input += tu.get("input_tokens", 0)
        total_output += tu.get("output_tokens", 0)
        total_cost += tu.get("estimated_cost_usd", 0.0)
        query_count += 1

    return {
        "query_count": query_count,
        "total_input_tokens": total_input,
        "total_output_tokens": total_output,
        "total_tokens": total_input + total_output,
        "total_cost_usd": round(total_cost, 4),
    }


@router.post("/", response_model=QueryResponse)
async def query_endpoint(
    request: Request,
    body: QueryRequest,
    background_tasks: BackgroundTasks,
    user: dict = Depends(get_current_user_or_api_key),
):
    # Sanitize input
    query = body.query.strip()
    if not query:
        raise HTTPException(status_code=400, detail="Query cannot be empty")

    # Free-trial gate (owners/paid bypass)
    _enforce_free_limit(user)

    # Scope session to user
    scoped_session = f"{user['user_id']}:{body.session_id}"

    # Use org_id as namespace so queries search the user's own documents
    namespace = user.get("org_id") or "default"

    try:
        response, analysis = ask(
            query, scoped_session,
            namespace=namespace,
            language=body.language,
            source_filters=body.source_filters or [],
            user_id=user["user_id"],
            force_bi=body.force_bi,
            model_override=body.model or None,
            no_cache=body.no_cache,
        )
    except Exception as e:
        err_str = str(e).lower()
        logger.error("query_failed", error=str(e), user_id=user["user_id"])
        if "credit balance is too low" in err_str or "insufficient_quota" in err_str:
            raise HTTPException(
                status_code=402,
                detail="Anthropic API credits are exhausted. Please add credits at console.anthropic.com or switch to a Groq/Llama model.",
            )
        raise HTTPException(status_code=500, detail="Query processing failed")

    ip = request.headers.get("X-Forwarded-For", request.client.host if request.client else None)
    write_audit(
        event_type="query",
        user_id=user["user_id"],
        email=user.get("email"),
        org_id=user.get("org_id"),
        detail={
            "query": query[:500],
            "patterns": response.patterns_used,
            "quality_score": response.quality_score,
            "latency_ms": response.latency_ms,
            "channel": response.retrieval_channel,
            "cache_hit": response.verified_knowledge_hit,
            "token_usage": response.token_usage,
        },
        ip_address=ip,
    )

    logger.info(
        "query_complete",
        user_id=user["user_id"],
        patterns=response.patterns_used,
        quality=response.quality_score,
        latency_ms=response.latency_ms,
    )

    tu = response.token_usage or {}
    return QueryResponse(
        answer=response.answer_text,
        quality_score=response.quality_score,
        confidence=response.confidence,
        faithfulness=response.faithfulness,
        patterns_used=response.patterns_used,
        latency_ms=response.latency_ms,
        retrieval_channel=response.retrieval_channel,
        fallback_used=response.fallback_used,
        verified_knowledge_hit=response.verified_knowledge_hit,
        suggested_followups=response.suggested_followups,
        citation_map=response.citation_map,
        figures=getattr(response, "figures", []),
        session_id=body.session_id,
        # Report the model that ACTUALLY produced the answer (reflects fallbacks
        # such as Claude -> Groq on an Anthropic billing error), not just the
        # requested model — keeps the top badge consistent with token usage.
        model_used=(tu.get("model") or body.model),
        token_usage=TokenUsage(**tu) if tu else None,
    )


# ── Image-based batch query ────────────────────────────────────────────────────

ALLOWED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/gif", "image/bmp"}
MAX_IMAGE_BYTES = 10 * 1024 * 1024  # 10 MB


class ImageQueryItem(BaseModel):
    question: str
    answer: str
    quality_score: float
    confidence: str
    patterns_used: list[str]
    latency_ms: int
    citation_map: dict
    suggested_followups: list[str]


class ImageQueryResponse(BaseModel):
    questions_found: int
    results: list[ImageQueryItem]
    extraction_note: str = ""


def _extract_questions_from_image(image_bytes: bytes, content_type: str) -> list[str]:
    """Use Groq Vision to extract a numbered list of questions from an image."""
    ext = content_type.split("/")[-1].replace("jpeg", "jpg")
    mime_map = {
        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "webp": "image/webp", "gif": "image/gif", "bmp": "image/bmp",
    }
    mime = mime_map.get(ext, "image/png")
    b64 = base64.standard_b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"

    prompt = (
        "This image contains one or more questions (e.g. a question paper, exam sheet, or form).\n"
        "Extract every question exactly as written, one per numbered line (1. 2. 3. ...).\n"
        "IMPORTANT: If a question has multiple-choice options (a), (b), (c), (d) / A B C D, "
        "keep the question stem AND all its options together as ONE single numbered item "
        "(put the options on the same line, separated by ' | '). Do NOT output the options "
        "as separate questions.\n"
        "Do NOT include the marked/correct answer, instructions, or any other text.\n"
        "If no questions are found, output: NO_QUESTIONS_FOUND"
    )

    try:
        from groq import Groq
        client = Groq()
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            max_tokens=2048,
            messages=[{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": data_url}},
                {"type": "text", "text": prompt},
            ]}],
        )
        raw = response.choices[0].message.content.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Vision extraction failed: {e}")

    return _parse_question_lines(raw)


def _parse_question_lines(raw: str) -> list[str]:
    """Turn an LLM's numbered-list output into a clean list of questions.

    Defensive (LLM-independent): even if the model splits a multiple-choice
    question into separate lines, a standalone option line (e.g. "(a) ...",
    "b) ...", "(ii) ...") is merged back into the preceding question instead of
    becoming its own question. Answer-key lines ("Correct Answer: ...") are
    dropped. Works for any question type, any subject — no hardcoding.
    """
    if "NO_QUESTIONS_FOUND" in raw:
        return []
    import re
    # A line that is just a multiple-choice option (latin a-h or roman i/v/x).
    OPTION_RE = re.compile(r'^\(?(?:[a-hA-H]|[ivxIVX]{1,4})[\)\.]\s+\S')
    ANSWER_RE = re.compile(r'^\s*(?:correct\s+answer|answer|ans)\b\s*[:\-.]', re.IGNORECASE)
    questions: list[str] = []
    for line in raw.splitlines():
        line = line.strip()
        if not line:
            continue
        # Strip leading numbering like "1." "1)" "Q1."
        cleaned = re.sub(r'^(?:Q?\d+[\.\)]\s*)', '', line).strip()
        if not cleaned:
            continue
        # Drop standalone answer-key lines.
        if ANSWER_RE.match(cleaned):
            continue
        # Merge a stray standalone option into the previous question.
        if OPTION_RE.match(cleaned) and questions:
            questions[-1] = f"{questions[-1]} | {cleaned}"
            continue
        questions.append(cleaned)
    return questions


def _extract_text_from_upload(file_bytes: bytes, filename: str, content_type: str) -> str:
    """Extract plain text from any supported file (PDF, DOCX, TXT, MD, CSV)."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    ctype = content_type or ""

    if ext == ".pdf" or ctype == "application/pdf":
        import io
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes), strict=False)
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext in {".docx", ".doc"} or "word" in ctype or "officedocument" in ctype:
        import io
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # Plain text / markdown / csv / anything else — decode as text
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1", errors="replace")


def _extract_questions_from_text(text: str) -> list[str]:
    """Use a Groq text model to extract a numbered list of questions from text."""
    text = text.strip()
    if not text:
        return []
    snippet = text[:20000]  # keep the prompt bounded
    prompt = (
        "The following content contains one or more questions (e.g. a question paper, "
        "exam sheet, worksheet, or form).\n"
        "Extract every question exactly as written, one per numbered line (1. 2. 3. ...).\n"
        "IMPORTANT: If a question has multiple-choice options (a), (b), (c), (d) / A B C D, "
        "keep the question stem AND all its options together as ONE single numbered item "
        "(put the options on the same line, separated by ' | '). Do NOT output the options "
        "as separate questions.\n"
        "Do NOT include the marked/correct answer, instructions, or any other text.\n"
        "If no questions are found, output: NO_QUESTIONS_FOUND\n\n"
        f"CONTENT:\n{snippet}"
    )
    try:
        from groq import Groq
        client = Groq()
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            max_tokens=2048,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.choices[0].message.content.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Question extraction failed: {e}")
    return _parse_question_lines(raw)


@router.post("/from-image", response_model=ImageQueryResponse)
async def query_from_image(
    request: Request,
    file: UploadFile = File(...),
    language: str = Form(default="American English"),
    session_id: str = Form(default="default"),
    user: dict = Depends(get_current_user_or_api_key),
):
    _enforce_free_limit(user)
    file_bytes = await file.read()
    if len(file_bytes) > MAX_IMAGE_BYTES:
        raise HTTPException(status_code=400, detail="File too large — maximum 10 MB.")

    ctype = file.content_type or ""
    filename = file.filename or "upload"

    # Images → Groq Vision; any other file (PDF/DOCX/TXT/CSV/…) → text extraction.
    if ctype in ALLOWED_IMAGE_TYPES or ctype.startswith("image/"):
        questions = _extract_questions_from_image(file_bytes, ctype)
        source_label = "image"
    else:
        try:
            text = _extract_text_from_upload(file_bytes, filename, ctype)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Could not read file '{filename}': {e}")
        questions = _extract_questions_from_text(text)
        source_label = "file"

    if not questions:
        return ImageQueryResponse(questions_found=0, results=[], extraction_note=f"No questions detected in the uploaded {source_label}.")

    namespace = user.get("org_id") or "default"
    scoped_session = f"{user['user_id']}:{session_id}"

    results: list[ImageQueryItem] = []
    for q in questions:
        try:
            response, _ = ask(
                q, scoped_session,
                namespace=namespace,
                language=language,
                source_filters=[],
                user_id=user["user_id"],
            )
            results.append(ImageQueryItem(
                question=q,
                answer=response.answer_text,
                quality_score=response.quality_score,
                confidence=response.confidence,
                patterns_used=response.patterns_used,
                latency_ms=response.latency_ms,
                citation_map=response.citation_map,
                suggested_followups=response.suggested_followups,
            ))
        except Exception as e:
            results.append(ImageQueryItem(
                question=q,
                answer=f"Error processing this question: {e}",
                quality_score=0.0,
                confidence="LOW",
                patterns_used=[],
                latency_ms=0,
                citation_map={},
                suggested_followups=[],
            ))

    # Count this as one metered question toward the free-trial limit.
    write_audit(event_type="query", user_id=user["user_id"], org_id=user.get("org_id"),
                detail={"kind": "image_batch", "questions_found": len(questions)})

    return ImageQueryResponse(
        questions_found=len(questions),
        results=results,
        extraction_note=f"Extracted {len(questions)} question(s) from {source_label}.",
    )


# ── Answer Evaluation ─────────────────────────────────────────────────────────
# Upload a questions file + an answers file; for each question a reference answer
# is generated from the ingested documents (RAG), the student's answer is graded
# against it, scored out of 100, and a report with mistakes + corrections is
# produced. Works with any file type (PDF/DOCX/TXT/CSV/image) for both uploads.

_ANSWER_VISION_PROMPT = (
    "This image contains a student's written ANSWERS (e.g. an answer sheet).\n"
    "Extract each answer exactly as written, in order, one answer per line.\n"
    "Output ONLY the answers, numbered 1. 2. 3. etc. matching the answer order.\n"
    "Preserve the full text of each answer. If none are found, output: NO_QUESTIONS_FOUND"
)
_ANSWER_TEXT_PROMPT = (
    "The following content contains a student's written ANSWERS (an answer sheet).\n"
    "Extract each answer exactly as written, in order, one answer per line.\n"
    "Output ONLY the answers, numbered 1. 2. 3. etc. matching the answer order.\n"
    "Preserve the full text of each answer. If none are found, output: NO_QUESTIONS_FOUND"
)


def _extract_answers_from_text(text: str) -> list[str]:
    """Extract a numbered list of a student's answers from plain text."""
    text = text.strip()
    if not text:
        return []
    try:
        from groq import Groq
        client = Groq()
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            max_tokens=4096,
            messages=[{"role": "user", "content": f"{_ANSWER_TEXT_PROMPT}\n\nCONTENT:\n{text[:24000]}"}],
        )
        raw = response.choices[0].message.content.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Answer extraction failed: {e}")
    return _parse_question_lines(raw)


def _extract_answers_from_image(image_bytes: bytes, content_type: str) -> list[str]:
    """Extract a numbered list of a student's answers from an image (Groq Vision)."""
    ext = content_type.split("/")[-1].replace("jpeg", "jpg")
    mime_map = {"png": "image/png", "jpg": "image/jpeg", "webp": "image/webp",
                "gif": "image/gif", "bmp": "image/bmp"}
    mime = mime_map.get(ext, "image/png")
    b64 = base64.standard_b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"
    try:
        from groq import Groq
        client = Groq()
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            max_tokens=4096,
            messages=[{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": data_url}},
                {"type": "text", "text": _ANSWER_VISION_PROMPT},
            ]}],
        )
        raw = response.choices[0].message.content.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Answer vision extraction failed: {e}")
    return _parse_question_lines(raw)


def _extract_items_from_file(file_bytes: bytes, filename: str, content_type: str, kind: str) -> list[str]:
    """Extract a numbered list of questions or answers from any uploaded file."""
    ctype = content_type or ""
    is_image = ctype in ALLOWED_IMAGE_TYPES or ctype.startswith("image/")
    if kind == "answer":
        if is_image:
            return _extract_answers_from_image(file_bytes, ctype)
        return _extract_answers_from_text(_extract_text_from_upload(file_bytes, filename, ctype))
    # kind == "question"
    if is_image:
        return _extract_questions_from_image(file_bytes, ctype)
    return _extract_questions_from_text(_extract_text_from_upload(file_bytes, filename, ctype))


def _evaluate_answer(question: str, student_answer: str, reference_answer: str) -> dict:
    """Grade a student's answer against a reference answer; return a structured report."""
    import json
    prompt = (
        "You are a fair, question-aware examiner. Grade the STUDENT ANSWER out of 100.\n\n"
        "GOLDEN RULES:\n"
        "0. GRADE BY MEANING, NOT WORDING: Judge whether the student's answer is "
        "CONCEPTUALLY correct. Award full credit for correct answers written in the "
        "student's OWN WORDS, even if the phrasing, order, or vocabulary differ from the "
        "reference. NEVER penalise for not matching the reference's exact words, sentences, "
        "or structure. Reward paraphrases, synonyms, and equivalent explanations. Only the "
        "correctness of the meaning matters.\n"
        "1. QUESTION-AWARE: Grade ONLY against what the QUESTION explicitly asks. The "
        "REFERENCE may contain extra detail beyond the question — IGNORE any part the "
        "question did not request. NEVER deduct marks for omitting information that was "
        "not asked for (e.g. examples, applications, or implications the question never "
        "mentioned).\n"
        "2. REWARD CORRECT CONCEPTS: A correct core concept should already score high "
        "(~80+). Missing OPTIONAL detail is at most a small deduction, not a heavy one. "
        "Only require an example/derivation if the question literally asks for it.\n"
        "3. NUMERICAL QUESTIONS (asks to calculate/find a value): award partial credit by "
        "step — correct formula/approach 30, correct substitution 30, correct final value "
        "30, correct units 10. VERIFY THE ARITHMETIC YOURSELF, step by step, before "
        "judging. Do NOT trust the reference's number blindly — recompute independently; if "
        "the reference's value looks wrong, grade against your own correct computation. If "
        "the question asks for an AMOUNT/MAGNITUDE, do NOT deduct for a missing +/- sign.\n"
        "4. CONCEPTUAL QUESTIONS: core concept 40, correct formula/principle 20, "
        "explanation 15, example (ONLY if requested) 10, units/symbols 5, plus any numeric "
        "part 10.\n"
        "5. GROUND in the reference/source for facts; if the reference is insufficient, say "
        "so honestly and grade on correct domain knowledge — be consistent (never call it "
        "'insufficient context' yet still award full marks, or vice versa).\n\n"
        f"QUESTION:\n{question}\n\n"
        f"REFERENCE ANSWER (from the documents — may include extra, unrequested detail):\n{reference_answer}\n\n"
        f"STUDENT ANSWER:\n{student_answer}\n\n"
        "Return ONLY valid JSON with this exact shape:\n"
        '{"score": <int 0-100>, "verdict": "<correct|partially correct|incorrect>", '
        '"mistakes": ["<specific mistake actually relevant to the question>", ...], '
        '"corrections": ["<the correction for each mistake>", ...], '
        '"feedback": "<one short paragraph: the rubric-based justification of the score>"}\n'
        "If the student answer is missing or empty, score 0. Keep mistakes/corrections "
        "empty when the answer is fully correct. Do not list omissions that were not asked."
    )
    try:
        from groq import Groq
        client = Groq()
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            max_tokens=2048,
            temperature=0.0,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.choices[0].message.content.strip()
        start, end = raw.find("{"), raw.rfind("}") + 1
        data = json.loads(raw[start:end])
    except Exception as e:
        return {"score": 0, "verdict": "error", "mistakes": [f"Evaluation failed: {e}"],
                "corrections": [], "feedback": "Could not evaluate this answer."}
    try:
        data["score"] = max(0, min(100, int(round(float(data.get("score", 0))))))
    except Exception:
        data["score"] = 0
    data.setdefault("verdict", "")
    data.setdefault("mistakes", [])
    data.setdefault("corrections", [])
    data.setdefault("feedback", "")
    return data


class EvalItem(BaseModel):
    question: str
    student_answer: str
    reference_answer: str
    score: int
    verdict: str
    mistakes: list[str]
    corrections: list[str]
    feedback: str


class EvalResponse(BaseModel):
    overall_score: int
    total_questions: int
    results: list[EvalItem]
    note: str = ""


@router.post("/evaluate", response_model=EvalResponse)
async def evaluate_answers(
    request: Request,
    questions_file: UploadFile = File(...),
    answers_file: UploadFile = File(...),
    language: str = Form(default="American English"),
    session_id: str = Form(default="default"),
    source_filters: str = Form(default=""),
    user: dict = Depends(get_current_user_or_api_key),
):
    _enforce_free_limit(user)
    q_bytes = await questions_file.read()
    a_bytes = await answers_file.read()
    if len(q_bytes) > MAX_IMAGE_BYTES or len(a_bytes) > MAX_IMAGE_BYTES:
        raise HTTPException(status_code=400, detail="File too large — maximum 10 MB each.")

    # source_filters arrives as a JSON array string; scope retrieval to those
    # ingested files so evaluation only checks against the chosen documents.
    filters: list[str] = []
    if source_filters.strip():
        import json as _json
        try:
            parsed = _json.loads(source_filters)
            if isinstance(parsed, list):
                filters = [str(s) for s in parsed if str(s).strip()]
        except Exception:
            filters = [s.strip() for s in source_filters.split(",") if s.strip()]

    questions = _extract_items_from_file(
        q_bytes, questions_file.filename or "questions", questions_file.content_type or "", "question")
    answers = _extract_items_from_file(
        a_bytes, answers_file.filename or "answers", answers_file.content_type or "", "answer")

    if not questions:
        return EvalResponse(overall_score=0, total_questions=0, results=[],
                            note="No questions detected in the questions file.")

    namespace = user.get("org_id") or "default"
    scoped_session = f"{user['user_id']}:{session_id}"

    results: list[EvalItem] = []
    for i, q in enumerate(questions):
        student_answer = answers[i] if i < len(answers) else ""
        # Reference answer from the ingested documents (RAG)
        try:
            ref_resp, _ = ask(q, scoped_session, namespace=namespace, language=language,
                              source_filters=filters, user_id=user["user_id"])
            reference = ref_resp.answer_text
        except Exception as e:
            reference = f"(Could not generate reference answer: {e})"
        ev = _evaluate_answer(q, student_answer, reference)
        results.append(EvalItem(
            question=q,
            student_answer=student_answer,
            reference_answer=reference,
            score=ev["score"],
            verdict=ev["verdict"],
            mistakes=ev["mistakes"],
            corrections=ev["corrections"],
            feedback=ev["feedback"],
        ))

    overall = round(sum(r.score for r in results) / len(results)) if results else 0
    note = f"Evaluated {len(results)} answer(s)."
    if len(answers) < len(questions):
        note += f" {len(questions) - len(answers)} question(s) had no matching answer (scored 0)."
    # Count this evaluation run as one metered question toward the free-trial limit.
    write_audit(event_type="query", user_id=user["user_id"], org_id=user.get("org_id"),
                detail={"kind": "answer_evaluation", "questions": len(questions)})
    return EvalResponse(
        overall_score=overall,
        total_questions=len(questions),
        results=results,
        note=note,
    )


# ── Podcast audio (Text-to-Speech) ────────────────────────────────────────────
# Generate a downloadable MP3 of the Q&A using Google Cloud Text-to-Speech via
# REST (SA token — no extra dependency). The texttospeech API is already enabled.

def _tts_voice_for_language(language: str) -> tuple[str, str]:
    """Map an answer language to a (languageCode, voiceName) for Cloud TTS."""
    lang = (language or "").lower()
    if "tamil" in lang:
        return ("ta-IN", "ta-IN-Standard-A")
    if "indian" in lang:
        return ("en-IN", "en-IN-Standard-A")
    if "british" in lang:
        return ("en-GB", "en-GB-Standard-A")
    if "australian" in lang:
        return ("en-AU", "en-AU-Standard-A")
    if "hindi" in lang:
        return ("hi-IN", "hi-IN-Standard-A")
    return ("en-US", "en-US-Standard-C")


def _chunk_text_for_tts(text: str, limit: int = 4500) -> list[str]:
    """Split text into <=limit-char pieces on sentence boundaries (TTS caps at 5000 bytes)."""
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks: list[str] = []
    cur = ""
    for s in sentences:
        if len(cur) + len(s) + 1 > limit:
            if cur:
                chunks.append(cur)
            # a single very long sentence — hard-split it
            while len(s) > limit:
                chunks.append(s[:limit])
                s = s[limit:]
            cur = s
        else:
            cur = f"{cur} {s}".strip()
    if cur:
        chunks.append(cur)
    return chunks


def _synthesize_tts(text: str, language: str = "American English") -> bytes:
    """Synthesize MP3 audio from text via Cloud TTS REST. Concatenates chunks."""
    import json as _json
    import urllib.request
    import google.auth
    from google.auth.transport import requests as ga_requests

    creds, _ = google.auth.default(scopes=["https://www.googleapis.com/auth/cloud-platform"])
    creds.refresh(ga_requests.Request())
    lang_code, voice = _tts_voice_for_language(language)

    audio = b""
    for chunk in _chunk_text_for_tts(text):
        body = {
            "input": {"text": chunk},
            "voice": {"languageCode": lang_code, "name": voice},
            "audioConfig": {"audioEncoding": "MP3"},
        }
        req = urllib.request.Request(
            "https://texttospeech.googleapis.com/v1/text:synthesize",
            data=_json.dumps(body).encode(), method="POST",
            headers={"Authorization": f"Bearer {creds.token}", "Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=60) as r:  # nosec B310
            resp = _json.loads(r.read())
        audio += base64.b64decode(resp["audioContent"])
    return audio


class TTSRequest(BaseModel):
    text: str = Field(..., min_length=1, max_length=20000)
    language: str = Field(default="American English", max_length=50)


@router.post("/tts")
async def tts_endpoint(body: TTSRequest, user: dict = Depends(get_current_user_or_api_key)):
    """Return a downloadable MP3 (podcast) of the supplied Q&A text."""
    text = body.text.strip()
    if not text:
        raise HTTPException(status_code=400, detail="No text to synthesize.")
    try:
        audio = _synthesize_tts(text, body.language)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Audio generation failed: {e}")
    return Response(
        content=audio,
        media_type="audio/mpeg",
        headers={"Content-Disposition": 'attachment; filename="maximai-answer.mp3"'},
    )


# ── Downloads: evaluation PDF + answer slides (PPTX) ──────────────────────────

def _latin1(s: str) -> str:
    """fpdf2 core fonts are latin-1; replace unsupported chars so export never fails."""
    return (s or "").encode("latin-1", "replace").decode("latin-1")


def _build_eval_pdf(data: "EvalResponse") -> bytes:
    # Explicit width (effective page width) + cursor control — robust on fpdf2 2.8+
    # (where the deprecated `ln=` was removed and bare multi_cell(0,...) can misbehave).
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    W = pdf.epw  # effective page width (page minus margins)

    def line(text: str, h: float = 6, style: str = "", size: int = 10):
        pdf.set_font("Helvetica", style, size)
        pdf.multi_cell(W, h, _latin1(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    line("MaximAI - Answer Evaluation Report", h=10, style="B", size=16)
    line(f"Overall Score: {data.overall_score}/100   |   {data.total_questions} question(s)", h=8, style="B", size=12)
    if data.note:
        line(data.note, style="I")
    pdf.ln(2)

    for i, r in enumerate(data.results, 1):
        line(f"Q{i}: {r.question}", h=7, style="B", size=12)
        line(f"Score: {r.score}/100  ({r.verdict})", h=7, style="B", size=11)
        line(f"Student answer: {r.student_answer or '(no answer)'}")
        if r.mistakes:
            line("Mistakes:", style="B")
            for m in r.mistakes:
                line(f"  - {m}")
        if r.corrections:
            line("Corrections:", style="B")
            for c in r.corrections:
                line(f"  - {c}")
        if r.feedback:
            line(f"Feedback: {r.feedback}", style="I")
        pdf.ln(4)
    return bytes(pdf.output())


@router.post("/evaluate/pdf")
async def evaluate_pdf(body: EvalResponse, user: dict = Depends(get_current_user_or_api_key)):
    """Return the evaluation results as a single downloadable PDF."""
    try:
        data = _build_eval_pdf(body)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {e}")
    return Response(content=data, media_type="application/pdf",
                    headers={"Content-Disposition": 'attachment; filename="maximai-evaluation.pdf"'})


class SlidesRequest(BaseModel):
    question: str = Field(..., min_length=1, max_length=4000)
    answer: str = Field(..., min_length=1, max_length=40000)


def _build_slides_pptx(question: str, answer: str) -> bytes:
    import io
    import re
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "MaximAI — Answer"
    try:
        s.placeholders[1].text = question[:250]
    except Exception:
        pass

    # Plain-text lines from the (markdown) answer; '#' headings start new slides.
    raw_lines = [ln.rstrip() for ln in answer.splitlines()]
    sections: list[tuple[str, list[str]]] = []
    cur_title, cur_body = "Answer", []
    for ln in raw_lines:
        h = re.match(r'^\s{0,3}#{1,6}\s+(.*)', ln)
        if h:
            if cur_body:
                sections.append((cur_title, cur_body))
            cur_title, cur_body = h.group(1).strip()[:80] or "Answer", []
        else:
            t = re.sub(r'[*`>]', '', ln).strip()
            if t:
                cur_body.append(t)
    if cur_body:
        sections.append((cur_title, cur_body))
    if not sections:
        sections = [("Answer", [re.sub(r'[*`>#]', '', answer).strip()[:1500] or " "])]

    # Each section -> one or more content slides (cap bullets/slide to fit).
    for title, body in sections:
        for j in range(0, len(body), 10):
            chunk = body[j:j + 10]
            cs = prs.slides.add_slide(prs.slide_layouts[1])
            cs.shapes.title.text = title if j == 0 else f"{title} (cont.)"
            tf = cs.placeholders[1].text_frame
            tf.text = chunk[0][:200]
            for line in chunk[1:]:
                p = tf.add_paragraph()
                p.text = line[:200]
                p.font.size = Pt(16)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@router.post("/slides")
async def answer_slides(body: SlidesRequest, user: dict = Depends(get_current_user_or_api_key)):
    """Return a PowerPoint (.pptx) of the question + answer."""
    try:
        data = _build_slides_pptx(body.question, body.answer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Slide generation failed: {e}")
    return Response(
        content=data,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="maximai-answer.pptx"'},
    )
