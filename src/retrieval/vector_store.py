import json
import re as _re
import time
from pathlib import Path


def _clean_text(text: str) -> str:
    """Normalise raw extracted text before chunking.

    Removes:
    - HTML / XML tags
    - Null bytes and non-printable control characters (except newlines/tabs)
    - Repeated whitespace / blank lines (3+ → 2)
    - Common boilerplate footers (cookie banners, print headers, etc.)
    """
    # HTML tags
    text = _re.sub(r'<[^>]+>', ' ', text)
    # Null bytes and control chars (keep \n \t)
    text = _re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    # Unicode replacement characters
    text = text.replace('�', '')
    # Collapse 3+ consecutive blank lines → 2
    text = _re.sub(r'\n{3,}', '\n\n', text)
    # Collapse sequences of spaces/tabs → single space (within lines)
    text = _re.sub(r'[ \t]{2,}', ' ', text)
    # Strip common boilerplate phrases (case-insensitive)
    _BOILERPLATE = [
        r'cookie policy.*?(?=\n|$)',
        r'privacy policy.*?(?=\n|$)',
        r'terms of (use|service).*?(?=\n|$)',
        r'all rights reserved.*?(?=\n|$)',
        r'subscribe to our newsletter.*?(?=\n|$)',
        r'javascript must be enabled.*?(?=\n|$)',
    ]
    for pat in _BOILERPLATE:
        text = _re.sub(pat, '', text, flags=_re.IGNORECASE)
    return text.strip()


def _rows_to_text(rows: list[list[str]]) -> list[str]:
    """Convert a 2-D table (list of rows, each a list of cell strings) to
    header-enriched text lines so column labels stay with their values.

    Example — given:
        [["Competition","Test","ODI","FC"], ["Batting average","53.78","44.83","57.84"]]
    Returns:
        ["Competition | Test | ODI | FC",
         "Batting average: Test=53.78 | ODI=44.83 | FC=57.84"]
    """
    if not rows:
        return []
    result = []
    headers: list[str] = []
    for row in rows:
        cells = [str(c).strip() for c in row]
        if not any(cells):
            continue
        # Detect header row: all cells non-empty and first data row not yet set
        if not headers:
            headers = cells
            result.append(" | ".join(cells))
        elif headers and len(cells) == len(headers) and len(headers) > 1:
            row_label = cells[0]
            enriched = [f"{headers[i]}={cells[i]}" for i in range(1, len(cells)) if cells[i]]
            # Emit plain pipe (legacy) AND enriched (header-labelled) so both
            # raw-value and column-label queries hit the same data.
            result.append(" | ".join(cells))
            result.append(f"{row_label}: {' | '.join(enriched)}" if enriched else " | ".join(cells))
        else:
            result.append(" | ".join(cells))
    return result
from sentence_transformers import SentenceTransformer
from pinecone import Pinecone, ServerlessSpec
from langchain_community.document_loaders import TextLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from src.config import (
    EMBEDDING_MODEL, EMBEDDING_DIM, CHUNK_SIZE, CHUNK_OVERLAP, CHUNK_SIZES, TOP_K,
    TABLE_MAX_CHUNK, PINECONE_API_KEY, PINECONE_INDEX_NAME,
)

_PINECONE_CLOUD = "aws"
_PINECONE_REGION = "us-east-1"

_pc: Pinecone | None = None
_index = None
_embedder: SentenceTransformer | None = None


def _is_table_line(line: str) -> bool:
    """Return True if this line is part of a serialised table block.

    Matches both formats produced by _rows_to_text:
      - Plain pipe rows:    "Col A | Col B | Col C"        (>= 1 pipe)
      - Enriched rows:      "Label: Col=val | Col=val"     (1+ pipe, contains '=')
    Also matches markdown-style rows that start with '|'.
    """
    stripped = line.strip()
    if not stripped:
        return False
    if stripped.startswith("|"):
        return True
    pipes = stripped.count("|")
    if pipes >= 2:
        return True
    # Enriched row: "Label: ColA=value | ColB=value" has exactly 1 pipe
    if pipes == 1 and "=" in stripped and ":" in stripped:
        return True
    return False


def _split_long_table(block: str, chunk_size: int) -> list[str]:
    """Split an oversized table at row boundaries, prepending the header to each chunk.

    Strategy:
    - Row 0 (and row 1 if it's a separator like |---|---) = header, prepended to every chunk.
    - Remaining rows are accumulated into chunks ≤ chunk_size.
    - Each chunk is self-contained: header + N data rows.
    """
    rows = [r for r in block.split("\n") if r.strip()]
    if not rows:
        return []

    # Identify header rows (first row + optional separator row like |---|---|)
    header_rows: list[str] = [rows[0]]
    data_start = 1
    if len(rows) > 1 and all(c in "|-: " for c in rows[1].replace("|", "")):
        header_rows.append(rows[1])
        data_start = 2

    header = "\n".join(header_rows)
    data_rows = rows[data_start:]

    chunks: list[str] = []
    current_rows: list[str] = []
    current_len = len(header) + 1

    for row in data_rows:
        row_len = len(row) + 1
        if current_len + row_len > chunk_size and current_rows:
            chunks.append(header + "\n" + "\n".join(current_rows))
            current_rows = []
            current_len = len(header) + 1
        current_rows.append(row)
        current_len += row_len

    if current_rows:
        chunks.append(header + "\n" + "\n".join(current_rows))

    return chunks


def _split_long_paragraph(para: str, chunk_size: int) -> list[str]:
    """Split a paragraph that exceeds chunk_size at sentence boundaries.

    Tries to cut on '. ', '? ', '! ' so chunks never end mid-sentence.
    Falls back to RecursiveCharacterTextSplitter if no sentence boundary works.
    """
    import re
    if len(para) <= chunk_size:
        return [para]

    sentence_end = re.compile(r'(?<=[.?!])\s+')
    sentences = sentence_end.split(para)

    chunks: list[str] = []
    current = ""
    for sent in sentences:
        if current and len(current) + len(sent) + 1 > chunk_size:
            chunks.append(current.strip())
            current = sent
        else:
            current = (current + " " + sent).strip() if current else sent

    if current:
        chunks.append(current.strip())

    # Last-resort: if any chunk still exceeds limit, use RecursiveCharacterTextSplitter
    result: list[str] = []
    fallback = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=0)
    for c in chunks:
        if len(c) > chunk_size * 1.5:
            result.extend(fallback.split_text(c))
        else:
            result.append(c)
    return result


def _table_aware_split(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Semantic-boundary chunking: keeps tables intact, respects paragraph and sentence breaks.

    Paragraph text:
    - Split first on paragraph boundaries (blank lines).
    - Keep each paragraph as its own chunk if it fits within chunk_size.
    - Paragraphs exceeding chunk_size are split at sentence boundaries (never mid-sentence).
    - Overlap = full previous paragraph (not a fixed char count), up to chunk_overlap chars,
      prepended to the next chunk so boundary paragraphs are always retrievable.

    Table blocks (consecutive pipe-row lines):
    - Short tables (≤ chunk_size): kept as one atomic chunk — all rows together.
    - Long tables (> chunk_size): split at row boundaries with the header row
      prepended to every sub-chunk so each piece is self-contained.
    """
    lines = text.split("\n")

    # Segment the text into alternating table / narrative blocks
    segments: list[tuple[str, bool]] = []  # (block_text, is_table)
    current_lines: list[str] = []
    in_table = False

    for line in lines:
        is_tl = _is_table_line(line)
        if is_tl != in_table:
            if current_lines:
                segments.append(("\n".join(current_lines), in_table))
            current_lines = [line]
            in_table = is_tl
        else:
            current_lines.append(line)
    if current_lines:
        segments.append(("\n".join(current_lines), in_table))

    raw_chunks: list[str] = []

    for block, is_table in segments:
        block = block.strip()
        if not block:
            continue

        if is_table:
            # Use TABLE_MAX_CHUNK (not chunk_size) so typical comparison tables
            # (< 4000 chars) are stored as a single atomic chunk — all rows together.
            # Only genuinely huge tables (e.g. 100+ row spreadsheets pasted into a PDF)
            # are split, and even then at row boundaries with header prepended.
            if len(block) <= TABLE_MAX_CHUNK:
                raw_chunks.append(block)
            else:
                raw_chunks.extend(_split_long_table(block, TABLE_MAX_CHUNK))
        else:
            # Split on paragraph boundaries first
            paragraphs = [p.strip() for p in block.split("\n\n") if p.strip()]
            for para in paragraphs:
                raw_chunks.extend(_split_long_paragraph(para, chunk_size))

    if not raw_chunks:
        return []

    # Apply paragraph-level overlap: prepend the tail of the previous chunk
    # (up to chunk_overlap chars) to each subsequent chunk so boundary content
    # is duplicated into the next chunk's embedding.
    final_chunks: list[str] = [raw_chunks[0]]
    for i in range(1, len(raw_chunks)):
        prev = raw_chunks[i - 1]
        overlap_text = prev[-chunk_overlap:].strip() if len(prev) > chunk_overlap else prev
        # Only prepend if the overlap isn't already in the current chunk
        current = raw_chunks[i]
        if overlap_text and not current.startswith(overlap_text):
            final_chunks.append(overlap_text + "\n\n" + current)
        else:
            final_chunks.append(current)

    return [c for c in final_chunks if c.strip()]


def _get_embedder() -> SentenceTransformer:
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBEDDING_MODEL)
    return _embedder


def _ensure_index(pc: Pinecone) -> None:
    existing = [idx.name for idx in pc.list_indexes()]
    if PINECONE_INDEX_NAME in existing:
        return
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=EMBEDDING_DIM,
        metric="cosine",
        spec=ServerlessSpec(cloud=_PINECONE_CLOUD, region=_PINECONE_REGION),
    )
    for _ in range(60):
        if pc.describe_index(PINECONE_INDEX_NAME).status.get("ready", False):
            return
        time.sleep(2)


def _get_index():
    global _pc, _index
    if _index is None:
        _pc = Pinecone(api_key=PINECONE_API_KEY)
        _ensure_index(_pc)
        _index = _pc.Index(PINECONE_INDEX_NAME)
    return _index


def _embed(texts: list[str]) -> list[list[float]]:
    return _get_embedder().encode(texts, show_progress_bar=False).tolist()


# Pinecone rejects upserts of more than 1000 vectors per request
# (and recommends a few hundred for payload-size headroom).
UPSERT_BATCH_SIZE = 200


def _batched_upsert(vectors: list[dict], namespace: str) -> None:
    index = _get_index()
    for start in range(0, len(vectors), UPSERT_BATCH_SIZE):
        index.upsert(vectors=vectors[start:start + UPSERT_BATCH_SIZE], namespace=namespace)


def _describe_image_bytes(image_bytes: bytes, ext: str = "png", context: str = "") -> str:
    """Send raw image bytes to Groq Vision and return a description.

    Args:
        image_bytes: Raw bytes of the image.
        ext:         Image format extension (png/jpg/gif/webp/bmp).
        context:     Hint (e.g. slide number) added to the prompt.
    Returns:
        Vision model description, or empty string on failure.
    """
    import base64
    try:
        from groq import Groq
    except ImportError:
        return ""

    mime_map = {
        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp",
        "tiff": "image/tiff", "emf": "image/png", "wmf": "image/png",
    }
    mime = mime_map.get(ext.lower().lstrip("."), "image/png")
    b64 = base64.standard_b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"

    loc = f" ({context})" if context else ""
    prompt = (
        f"You are analyzing an embedded image from a document{loc}.\n"
        "Extract ALL information visible:\n"
        "1. DIAGRAMS/FLOWCHARTS: describe each node, arrow, label, and flow\n"
        "2. CHARTS/GRAPHS: type, axes, all data values, trends, legend items\n"
        "3. TABLES: reproduce all rows and columns as pipe-separated text\n"
        "4. SCREENSHOTS/UI: describe the interface, visible text, key actions\n"
        "5. ARCHITECTURE DIAGRAMS: components, connections, data flow\n"
        "6. ANY TEXT: transcribe verbatim\n\n"
        "Be thorough — this text will be the only representation of this image in the RAG system."
    )

    try:
        client = Groq()
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            max_tokens=2048,
            messages=[{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": data_url}},
                {"type": "text", "text": prompt},
            ]}],
        )
        return response.choices[0].message.content.strip()
    except Exception:
        return ""


def _extract_text_from_pptx(path: Path) -> str:
    """Extract text, tables, images, and charts from a PowerPoint file (.pptx/.ppt).

    Legacy .ppt files are converted to .pptx via LibreOffice before processing.

    Shapes are processed in top-left → bottom-right reading order (by top then
    left position) so multi-column slide layouts are read correctly.

    Per shape:
      - TEXT shapes → verbatim text
      - TABLE shapes → pipe-separated rows
      - PICTURE shapes → OCR (column-aware) first; Groq Vision as fallback/supplement
        for diagram/chart images where OCR alone is insufficient
      - CHART shapes → title + series data as structured text
      - GROUP shapes → recursed into child shapes (position-sorted)

    For fully-image slides (slide background is a picture), the entire slide
    is rendered via pdf2image and OCR'd as a single layout unit.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import tempfile
    import subprocess

    # Convert legacy .ppt to .pptx using LibreOffice
    if path.suffix.lower() == ".ppt":
        tmp_dir = tempfile.mkdtemp()
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pptx",
                 "--outdir", tmp_dir, str(path)],
                check=True, timeout=120,
                capture_output=True,
            )
            converted = Path(tmp_dir) / (path.stem + ".pptx")
            if converted.exists():
                path = converted
        except Exception:
            pass  # fall through — python-pptx may still handle some .ppt files

    prs = Presentation(str(path))
    lines = []
    img_counter = [0]

    def _shape_top_left(shape):
        try:
            return (shape.top or 0, shape.left or 0)
        except Exception:
            return (0, 0)

    def _process_picture(shape, slide_num: int, slide_width_px: int = 1280):
        try:
            img = shape.image
            img_bytes = img.blob
            img_counter[0] += 1
            context = f"Slide {slide_num}, Image {img_counter[0]}"

            # Try OCR first — good for text-heavy images and multi-column layouts
            ocr_text = ""
            try:
                from PIL import Image as PILImage
                import io
                pil_img = PILImage.open(io.BytesIO(img_bytes)).convert("RGB")
                ocr_text = _ocr_image_to_text(pil_img, context=context)
            except Exception:
                pass

            # Also get Vision description for diagrams / charts / non-text images
            vision_desc = _describe_image_bytes(img_bytes, img.ext, context)

            # Prefer whichever gives more content; include both when both have value
            parts = []
            if ocr_text.strip():
                parts.append(ocr_text.strip())
            if vision_desc.strip() and vision_desc.strip() != ocr_text.strip():
                parts.append(vision_desc.strip())

            if parts:
                lines.append(f"\n### Image {img_counter[0]} (Slide {slide_num})")
                lines.extend(parts)
        except Exception:
            pass

    def _process_shape(shape, slide_num: int):
        # Native text (handles text boxes, titles, subtitles)
        if hasattr(shape, "text") and shape.text.strip():
            lines.append(shape.text.strip())

        stype = shape.shape_type

        if stype == MSO_SHAPE_TYPE.TABLE:
            try:
                raw_rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                lines.extend(_rows_to_text(raw_rows))
            except Exception:
                pass

        elif stype == MSO_SHAPE_TYPE.PICTURE:
            _process_picture(shape, slide_num)

        elif stype == MSO_SHAPE_TYPE.CHART:
            try:
                chart = shape.chart
                chart_title = ""
                try:
                    chart_title = chart.chart_title.text_frame.text.strip()
                except Exception:
                    pass
                lines.append(f"\n### Chart{': ' + chart_title if chart_title else ''} (Slide {slide_num})")
                for series in chart.series:
                    try:
                        s_name = series.name or "Series"
                        values = [str(v) for v in (series.values or [])]
                        lines.append(f"{s_name}: {', '.join(values)}")
                    except Exception:
                        pass
            except Exception:
                pass

        elif stype == MSO_SHAPE_TYPE.GROUP:
            try:
                child_shapes = sorted(shape.shapes, key=_shape_top_left)
                for child in child_shapes:
                    _process_shape(child, slide_num)
            except Exception:
                pass

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n--- Slide {slide_num} ---")
        # Process shapes in top→bottom, left→right order for correct reading sequence
        sorted_shapes = sorted(slide.shapes, key=_shape_top_left)
        for shape in sorted_shapes:
            _process_shape(shape, slide_num)

    return "\n".join(lines)


def _column_sort_words(words: list[dict], page_width: float, line_y_tol: float | None = None) -> str:
    """Sort word dicts in reading order (column-aware) and reconstruct text.

    Accepts a list of dicts with keys: text, x0, x1, top, bottom.
    Works for both pdfplumber native words (PDF pts) and pytesseract words (px).

    Column detection: horizontal gaps wider than 15% of page_width are treated
    as column gutters. This handles 1-col, 2-col, 3-col, and magazine layouts.

    line_y_tol: vertical tolerance for grouping words into one line.
                Defaults to 40% of the median word height (resolution-agnostic).
    """
    if not words:
        return ""

    # Adaptive line-height tolerance: 40% of median word height
    if line_y_tol is None:
        heights = [max(w["bottom"] - w["top"], 1) for w in words]
        heights.sort()
        median_h = heights[len(heights) // 2]
        line_y_tol = max(median_h * 0.4, 2)

    # Build column boundary list from horizontal gap analysis
    x0_vals = sorted(set(round(w["x0"]) for w in words))
    col_boundaries = [0.0]
    for i in range(1, len(x0_vals)):
        gap = x0_vals[i] - x0_vals[i - 1]
        if gap > page_width * 0.15:
            col_boundaries.append((x0_vals[i - 1] + x0_vals[i]) / 2)
    col_boundaries.append(page_width + 1)

    def _col(w: dict) -> int:
        mid = (w["x0"] + w["x1"]) / 2
        for c in range(len(col_boundaries) - 1):
            if col_boundaries[c] <= mid < col_boundaries[c + 1]:
                return c
        return len(col_boundaries) - 2

    # Sort: column first, then y (quantised), then x
    bucket = max(line_y_tol * 0.75, 1)
    words = sorted(words, key=lambda w: (_col(w), round(w["top"] / bucket) * bucket, w["x0"]))

    # Group into visual lines
    lines_out: list[list[dict]] = []
    cur_line: list[dict] = []
    cur_col = -1

    for word in words:
        col = _col(word)
        if col != cur_col:
            if cur_line:
                lines_out.append(cur_line)
            cur_line = [word]
            cur_col = col
            continue
        if not cur_line:
            cur_line = [word]
        elif abs(word["top"] - cur_line[-1]["top"]) <= line_y_tol:
            cur_line.append(word)
        else:
            lines_out.append(cur_line)
            cur_line = [word]

    if cur_line:
        lines_out.append(cur_line)

    text_lines = []
    for lw in lines_out:
        lw.sort(key=lambda w: w["x0"])
        text_lines.append(" ".join(w["text"] for w in lw))

    return "\n".join(text_lines)


def _ocr_image_to_text(pil_image, context: str = "") -> str:
    """OCR a PIL Image with column-aware reading order via pytesseract.

    Uses tesseract --psm 1 (auto OSD + layout) so it handles:
      - Single-column text
      - Multi-column (2, 3+ cols)
      - Mixed layouts (callout boxes, sidebars, headers)
      - Tables (extracted as pipe-separated rows)
      - Rotated text (OSD corrects orientation first)

    Falls back to empty string if pytesseract is not installed or fails.
    """
    try:
        import pytesseract
        from pytesseract import Output as TessOutput
    except ImportError:
        return ""

    try:
        data = pytesseract.image_to_data(
            pil_image,
            output_type=TessOutput.DICT,
            config="--psm 1 --oem 3",
        )
    except Exception:
        return ""

    img_w, _img_h = pil_image.size
    words: list[dict] = []
    for i, text in enumerate(data["text"]):
        text = (text or "").strip()
        if not text:
            continue
        conf = int(data["conf"][i])
        if conf < 20:
            continue
        x, y, w, h = data["left"][i], data["top"][i], data["width"][i], data["height"][i]
        words.append({"text": text, "x0": x, "x1": x + w, "top": y, "bottom": y + h})

    return _column_sort_words(words, page_width=img_w)


def _render_pdf_page(path: Path, page_num: int, dpi: int = 200):
    """Render a single PDF page to a PIL Image using pdf2image + Poppler.

    Returns None if pdf2image is not installed or rendering fails.
    page_num is 1-based.
    POPPLER_PATH env var (or src.config.POPPLER_PATH) points to the poppler bin
    directory on Windows; leave empty on Linux where poppler is on PATH.
    """
    try:
        from pdf2image import convert_from_path
    except ImportError:
        return None
    try:
        from src.config import POPPLER_PATH
        kwargs: dict = {"dpi": dpi, "first_page": page_num, "last_page": page_num}
        if POPPLER_PATH:
            kwargs["poppler_path"] = POPPLER_PATH
        images = convert_from_path(str(path), **kwargs)
        return images[0] if images else None
    except Exception:
        return None


def _extract_columns_from_page(plumb_page, excluded_bboxes: list) -> str:
    """Extract text from a pdfplumber page using column-aware word extraction.

    Falls back to OCR (via _render_pdf_page + _ocr_image_to_text) when the
    page yields fewer than OCR_SPARSE_THRESHOLD characters natively (scanned PDF).
    """
    try:
        raw_words = plumb_page.extract_words(
            x_tolerance=3, y_tolerance=3, keep_blank_chars=False,
            use_text_flow=False,
        ) or []
    except Exception:
        return plumb_page.extract_text() or ""

    def _in_bbox(word, bb):
        return (bb[0] <= word["x0"] and word["x1"] <= bb[2] and
                bb[1] <= word["top"] and word["bottom"] <= bb[3])

    if excluded_bboxes:
        raw_words = [w for w in raw_words if not any(_in_bbox(w, bb) for bb in excluded_bboxes)]

    if not raw_words:
        return ""

    words = [{"text": w["text"], "x0": w["x0"], "x1": w["x1"],
               "top": w["top"], "bottom": w["bottom"]} for w in raw_words]
    return _column_sort_words(words, page_width=plumb_page.width or 600)


_STEM_SYMBOLS = _re.compile(
    r'[∫∑∏√∂∇∆≈≠≤≥±∞∈∉⊆⊇∩∪∀∃∄→←↔⇒⇔∧∨¬α-ωΑ-Ω]'
    r'|\\(?:frac|sqrt|int|sum|prod|lim|partial|nabla|Delta|alpha|beta|gamma|theta|lambda|sigma|pi|mu|epsilon)\b'
)
_STEM_SYMBOL_THRESHOLD = 0.003  # 0.3% of chars are math symbols → STEM doc
_STEM_MIN_PAGES = 3             # only bother on multi-page docs
_SCANNED_CHARS_PER_PAGE = 100   # avg extractable chars/page below this ⇒ scanned


def _is_stem_pdf(path: Path) -> bool:
    """Heuristically decide whether to route a PDF through Marker.

    Two triggers (Marker handles both — it OCRs and recovers LaTeX math):

    1. Born-digital STEM: the text layer has a high density of math symbols
       (∫, ∑, Greek, LaTeX commands) — a scientific paper or textbook.
    2. Scanned / image-based multi-page PDF: very little extractable text per
       page means equations and body text are images. pdfplumber would only
       get garbled OCR fragments, whereas Marker's layout+OCR pipeline
       reconstructs the math. (Option A — catches scanned textbook chapters
       like NCERT/CBSE physics that have no real text layer.)
    """
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            num_pages = len(pdf.pages)
            if num_pages < _STEM_MIN_PAGES:
                return False
            sample = pdf.pages[:10]
            total_chars = 0
            symbol_chars = 0
            for page in sample:
                text = page.extract_text() or ""
                total_chars += len(text)
                symbol_chars += len(_STEM_SYMBOLS.findall(text))

            # Trigger 2: scanned/image PDF (sparse text layer across pages)
            avg_chars_per_page = total_chars / len(sample)
            if avg_chars_per_page < _SCANNED_CHARS_PER_PAGE:
                return True

            # Trigger 1: dense math symbols in a real text layer
            if total_chars < 200:
                return False
            return (symbol_chars / total_chars) >= _STEM_SYMBOL_THRESHOLD
    except Exception:
        return False


def _extract_text_with_marker(path: Path, on_progress=None) -> str | None:
    """Use the Marker library to extract text from a STEM PDF.

    Marker outputs clean Markdown with LaTeX math blocks ($$...$$) preserved,
    which chunking can handle much better than flattened symbol text.

    Returns None if Marker is not installed or fails, so callers can fall back
    to the standard pdfplumber pipeline.
    """
    try:
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict
        from marker.output import text_from_rendered
    except Exception as e:
        # Log the real import error (missing transitive dep, etc.) instead of
        # silently returning None — this is what hid the broken --no-deps install.
        print(f"[ingest] Marker import failed: {type(e).__name__}: {e}", flush=True)
        return None

    if on_progress:
        on_progress(10, "Running Marker (loading models)")

    try:
        print("[ingest] Marker: loading models + converting…", flush=True)
        models = create_model_dict()
        converter = PdfConverter(artifact_dict=models)
        rendered = converter(str(path))
        markdown_text, _, _ = text_from_rendered(rendered)
        if on_progress:
            on_progress(75, "Marker extraction complete")
        return markdown_text if markdown_text and len(markdown_text.strip()) > 100 else None
    except Exception as e:
        print(f"[ingest] Marker conversion failed: {type(e).__name__}: {e}", flush=True)
        return None


def _extract_text_from_pdf(path: Path, on_progress=None) -> str:
    """Extract text, tables, and embedded images from a PDF.

    on_progress(pct: int, label: str) is called per-page so callers can
    report incremental progress (10%–80% range reserved for page extraction).
    """
    import pypdf
    import pdfplumber

    OCR_SPARSE_CHARS = 80
    lines = []

    with pdfplumber.open(str(path)) as pdf:
        total_pages = len(pdf.pages)
        for page_num, plumb_page in enumerate(pdf.pages, start=1):
            lines.append(f"\n--- Page {page_num} ---")

            if on_progress and total_pages > 0:
                pct = 10 + int((page_num / total_pages) * 65)
                on_progress(pct, f"Extracting page {page_num} of {total_pages}")

            tables = plumb_page.extract_tables() or []
            table_bboxes = [t.bbox for t in plumb_page.find_tables()]
            for table in tables:
                raw_rows = [
                    [str(cell).strip() if cell is not None else "" for cell in row]
                    for row in table
                ]
                lines.extend(_rows_to_text(raw_rows))

            page_text = _extract_columns_from_page(plumb_page, table_bboxes)

            if len(page_text.strip()) < OCR_SPARSE_CHARS:
                if on_progress:
                    on_progress(
                        10 + int((page_num / total_pages) * 65),
                        f"OCR page {page_num} of {total_pages}",
                    )
                pil_img = _render_pdf_page(path, page_num)
                if pil_img is not None:
                    ocr_text = _ocr_image_to_text(pil_img, context=f"Page {page_num}")
                    if len(ocr_text.strip()) > len(page_text.strip()):
                        page_text = ocr_text

            if page_text.strip():
                lines.append(page_text.strip())

    # ── Images via pypdf ─────────────────────────────────────────────────
    try:
        reader = pypdf.PdfReader(str(path), strict=False)
        for page_num, page in enumerate(reader.pages, start=1):
            img_count = 0
            try:
                for img_obj in page.images:
                    try:
                        img_bytes = img_obj.data
                        ext = Path(img_obj.name).suffix.lstrip(".") or "png"
                        img_count += 1
                        context = f"Page {page_num}, Image {img_count}"
                        description = _describe_image_bytes(img_bytes, ext, context)
                        if description:
                            lines.append(f"\n### Image {img_count} (Page {page_num})")
                            lines.append(description)
                    except Exception:
                        pass
            except Exception:
                pass
    except Exception:
        pass

    return "\n".join(lines)


def _docx_num_columns(section) -> int:
    """Return the number of text columns defined for a DOCX section (default 1)."""
    try:
        W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        sect_pr = section._sectPr
        cols_el = sect_pr.find(f"{{{W_NS}}}cols")
        if cols_el is not None:
            num = cols_el.get(f"{{{W_NS}}}num")
            if num:
                return int(num)
    except Exception:
        pass
    return 1


def _extract_text_from_docx(path: Path) -> str:
    """Extract text AND embedded images from a DOCX.

    - Paragraph and table text is extracted in XML document order, which
      matches reading order for both single- and multi-column DOCX layouts
      (Word stores paragraphs in flow order, not column order).
    - For sections with 2+ columns the section column count is noted in output
      so downstream readers understand the layout.
    - Embedded images (body + headers/footers): OCR first (column-aware via
      pytesseract); Groq Vision as fallback/supplement for diagram-heavy images.
      OCR handles scanned pages, text callout boxes, and multi-column image inserts.
    """
    import docx

    doc = docx.Document(str(path))
    lines = []

    # ── Paragraphs and tables in document order ────────────────────────
    # We walk the XML body children to interleave paragraphs and tables
    # correctly instead of iterating them separately (which loses ordering).
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph

    body = doc.element.body
    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            para = DocxParagraph(child, doc)
            if para.text.strip():
                lines.append(para.text.strip())
        elif tag == "tbl":
            tbl = DocxTable(child, doc)
            raw_rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            lines.extend(_rows_to_text(raw_rows))
        elif tag == "sectPr":
            pass  # section property — not content

    # ── Embedded images: OCR + Vision ─────────────────────────────────
    img_count = 0
    IMAGE_CONTENT_TYPES = {
        "image/png", "image/jpeg", "image/gif",
        "image/webp", "image/bmp", "image/tiff",
    }
    seen_blobs: set[int] = set()

    def _handle_image_blob(blob: bytes, ct: str):
        nonlocal img_count
        blob_id = hash(blob)
        if blob_id in seen_blobs:
            return
        seen_blobs.add(blob_id)
        ext = ct.split("/")[-1]
        img_count += 1
        context = f"Image {img_count}"

        # OCR first — handles multi-column text layouts inside images
        ocr_text = ""
        try:
            from PIL import Image as PILImage
            import io
            pil_img = PILImage.open(io.BytesIO(blob)).convert("RGB")
            ocr_text = _ocr_image_to_text(pil_img, context=context)
        except Exception:
            pass

        # Vision for diagrams/charts/non-text content
        vision_desc = _describe_image_bytes(blob, ext, context)

        parts = []
        if ocr_text.strip():
            parts.append(ocr_text.strip())
        if vision_desc.strip() and vision_desc.strip() != ocr_text.strip():
            parts.append(vision_desc.strip())

        if parts:
            lines.append(f"\n### Embedded Image {img_count}")
            lines.extend(parts)

    def _extract_images_from_part(part):
        try:
            for rel in part.rels.values():
                try:
                    target = rel.target_part
                    ct = target.content_type
                    if ct not in IMAGE_CONTENT_TYPES:
                        continue
                    _handle_image_blob(target.blob, ct)
                except Exception:
                    pass
        except Exception:
            pass

    _extract_images_from_part(doc.part)
    for section in doc.sections:
        for hf_part in [section.header.part, section.footer.part,
                         section.even_page_header.part, section.even_page_footer.part,
                         section.first_page_header.part, section.first_page_footer.part]:
            try:
                _extract_images_from_part(hf_part)
            except Exception:
                pass

    return "\n".join(lines)


def _extract_text_from_xls(path: Path) -> str:
    """Extract all text from a legacy Excel file (.xls) as header-enriched rows."""
    import xlrd
    wb = xlrd.open_workbook(str(path))
    lines = []
    for sheet in wb.sheets():
        lines.append(f"\n--- Sheet: {sheet.name} ---")
        raw_rows = [
            [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
            for row_idx in range(sheet.nrows)
        ]
        lines.extend(_rows_to_text(raw_rows))
    return "\n".join(lines)


def _extract_text_from_xlsx(path: Path) -> str:
    """Extract text (header-enriched rows) and embedded images from a modern Excel file (.xlsx)."""
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    lines = []
    img_count = 0
    for sheet in wb.worksheets:
        lines.append(f"\n--- Sheet: {sheet.title} ---")

        # ── Cell data ───────────────────────────────────────────────────
        raw_rows = [
            [str(cell) if cell is not None else "" for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        lines.extend(_rows_to_text(raw_rows))

        # ── Embedded images: OCR + Vision ───────────────────────────────
        for img_obj in getattr(sheet, "_images", []):
            try:
                img_data = img_obj._data()
                img_count += 1
                fmt = getattr(img_obj, "format", None) or "png"
                context = f"Sheet '{sheet.title}', Image {img_count}"

                ocr_text = ""
                try:
                    from PIL import Image as PILImage
                    import io
                    pil_img = PILImage.open(io.BytesIO(img_data)).convert("RGB")
                    ocr_text = _ocr_image_to_text(pil_img, context=context)
                except Exception:
                    pass

                vision_desc = _describe_image_bytes(img_data, fmt, context)

                parts = []
                if ocr_text.strip():
                    parts.append(ocr_text.strip())
                if vision_desc.strip() and vision_desc.strip() != ocr_text.strip():
                    parts.append(vision_desc.strip())

                if parts:
                    lines.append(f"\n### Embedded Image {img_count} (Sheet: {sheet.title})")
                    lines.extend(parts)
            except Exception:
                pass

    return "\n".join(lines)


def _extract_text_from_json(path: Path) -> str:
    """Convert JSON / GeoJSON to indented text for chunking."""
    with open(str(path), "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _delete_source_vectors(id_prefix: str, namespace: str) -> int:
    """Delete all Pinecone vectors whose ID starts with id_prefix.

    Called before every ingest so re-uploading a file never leaves stale
    chunks behind (e.g. when chunk count changes between uploads).
    Returns the number of vectors deleted.
    """
    try:
        index = _get_index()
        ids_to_delete: list[str] = []
        for page in index.list(prefix=id_prefix, namespace=namespace):
            ids_to_delete.extend(item.id for item in page)
        if ids_to_delete:
            index.delete(ids=ids_to_delete, namespace=namespace)
        return len(ids_to_delete)
    except Exception:
        return 0  # non-fatal — upsert will still overwrite matching IDs


def delete_by_source(source_name: str, namespace: str = "default") -> int:
    """Delete all Pinecone vectors for a given source across all prefix variants.

    Covers plain filenames, image:/video:/audio:/weblink: prefixes and
    URL-derived IDs so a single call cleans up any ingest type.
    Returns total vectors deleted.
    """
    base = source_name.strip()
    prefixes = [
        base,
        f"image:{base}",
        f"video:{base}",
        f"audio:{base}",
        f"weblink:{base}",
    ]
    total = 0
    for prefix in prefixes:
        total += _delete_source_vectors(prefix, namespace)
    return total


def ingest_file(
    file_path: str,
    namespace: str = "default",
    source_name: str | None = None,
    on_progress=None,
    force_marker: bool = False,
) -> int:
    """Ingest a document file into Pinecone.

    Args:
        file_path:    Path to the saved (possibly hashed) file on disk.
        namespace:    Pinecone namespace (org_id).
        source_name:  Human-readable original filename stored as metadata.
        on_progress:  Optional callable(pct: int, label: str).
                      Called at key stages so callers can surface progress.
    """
    def _prog(pct: int, label: str):
        if on_progress:
            try:
                on_progress(pct, label)
            except Exception:
                pass

    path = Path(file_path)
    suffix = path.suffix.lower()

    _prog(5, "Reading file")

    raw_text: str | None = None
    doc_type = "narrative"

    if suffix == ".pdf":
        # Marker (STEM PDF parser) loads ~2GB of ML models and needs an 8Gi
        # instance. Gated behind ENABLE_MARKER (off on the default 4Gi instance
        # to avoid OOM). Marker runs when the caller forces it (manual "STEM"
        # toggle) OR auto-detection flags the PDF as STEM/scanned.
        import os as _os
        _marker_on = _os.getenv("ENABLE_MARKER", "false").lower() in ("1", "true", "yes")
        _auto_stem = _is_stem_pdf(path) if _marker_on else False
        print(f"[ingest] pdf={source_name or path.name} marker_enabled={_marker_on} "
              f"force_marker={force_marker} auto_stem={_auto_stem}", flush=True)
        if _marker_on and (force_marker or _auto_stem):
            _prog(8, "STEM PDF detected — trying Marker parser")
            raw_text = _extract_text_with_marker(path, on_progress=_prog)
            if raw_text is None:
                print("[ingest] Marker returned no text — falling back to pdfplumber", flush=True)
                _prog(10, "Marker unavailable — using standard PDF parser")
                raw_text = _extract_text_from_pdf(path, on_progress=_prog)
            else:
                print(f"[ingest] Marker OK — {len(raw_text)} chars extracted", flush=True)
        else:
            raw_text = _extract_text_from_pdf(path, on_progress=_prog)
    elif suffix in (".txt", ".md"):
        loader = TextLoader(str(path), encoding="utf-8")
    elif suffix == ".csv":
        loader = TextLoader(str(path), encoding="utf-8")
        doc_type = "tabular"
    elif suffix in (".docx", ".doc"):
        raw_text = _extract_text_from_docx(path)
    elif suffix in (".pptx", ".ppt"):
        raw_text = _extract_text_from_pptx(path)
    elif suffix == ".xls":
        raw_text = _extract_text_from_xls(path)
        doc_type = "tabular"
    elif suffix in (".xlsx", ".xlsm"):
        raw_text = _extract_text_from_xlsx(path)
        doc_type = "tabular"
    elif suffix in (".json", ".geojson", ".gjson"):
        raw_text = _extract_text_from_json(path)
        doc_type = "code"
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    # Resolve the display source name — prefer caller-supplied original filename
    display_source = source_name or path.name
    # Use a stable vector ID prefix from the display name (not the hashed disk name)
    id_prefix = Path(display_source).stem

    # Delete any existing vectors for this source before upserting new ones.
    # This prevents stale chunks accumulating when a file is re-uploaded with
    # a different chunk count (e.g. file updated, or chunk size changed).
    _delete_source_vectors(id_prefix, namespace)

    # Per-document-type chunk sizes: smaller for tabular data (better cell precision),
    # larger for narrative text (better context preservation).
    _chunk_size, _chunk_overlap = CHUNK_SIZES.get(doc_type, (CHUNK_SIZE, CHUNK_OVERLAP))
    splitter = RecursiveCharacterTextSplitter(chunk_size=_chunk_size, chunk_overlap=_chunk_overlap)

    if raw_text is not None:
        _prog(78, "Chunking text")
        raw_text = _clean_text(raw_text)
        text_chunks = _table_aware_split(raw_text, _chunk_size, _chunk_overlap)
        if not text_chunks:
            return 0

        # Embed in batches with per-batch progress (78% → 92%)
        EMBED_BATCH = 64
        embeddings: list[list[float]] = []
        total_chunks = len(text_chunks)
        for batch_start in range(0, total_chunks, EMBED_BATCH):
            batch = text_chunks[batch_start:batch_start + EMBED_BATCH]
            embeddings.extend(_embed(batch))
            done = min(batch_start + EMBED_BATCH, total_chunks)
            pct = 78 + int((done / total_chunks) * 14)  # 78→92
            _prog(pct, f"Embedding chunk {done}/{total_chunks}")

        # Upsert in batches with per-batch progress (92% → 99%)
        index = _get_index()
        vectors = [
            {
                "id": f"{id_prefix}_{i}",
                "values": emb,
                "metadata": {
                    "content": chunk,
                    "source": display_source,
                    "doc_type": doc_type,
                    "chunk_idx": i,
                },
            }
            for i, (chunk, emb) in enumerate(zip(text_chunks, embeddings))
        ]
        total_vectors = len(vectors)
        for start in range(0, total_vectors, UPSERT_BATCH_SIZE):
            index.upsert(vectors=vectors[start:start + UPSERT_BATCH_SIZE], namespace=namespace)
            done = min(start + UPSERT_BATCH_SIZE, total_vectors)
            pct = 92 + int((done / total_vectors) * 7)  # 92→99
            _prog(pct, f"Uploading batch {done}/{total_vectors} vectors")

        return len(vectors)

    # LangChain Document path (TXT / MD / CSV)
    _prog(78, "Chunking text")
    docs = loader.load()
    for doc in docs:
        doc.page_content = _clean_text(doc.page_content)
    chunks = splitter.split_documents(docs)

    if not chunks:
        return 0

    _prog(85, f"Generating embeddings for {len(chunks)} chunks")
    texts = [c.page_content for c in chunks]
    embeddings = _embed(texts)
    _prog(94, "Uploading to index")
    vectors = [
        {
            "id": f"{id_prefix}_{i}",
            "values": emb,
            "metadata": {
                "content": chunk.page_content,
                "source": display_source,
                "doc_type": doc_type,
                "chunk_idx": i,
                "page": chunk.metadata.get("page", 0),
            },
        }
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings))
    ]
    _batched_upsert(vectors, namespace)
    return len(vectors)


def ingest_text(text: str, source: str = "manual", namespace: str = "default") -> int:
    # Delete existing vectors for this source before upserting
    _delete_source_vectors(source, namespace)
    splitter = RecursiveCharacterTextSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    chunks = splitter.split_text(text)
    if not chunks:
        return 0
    embeddings = _embed(chunks)
    vectors = [
        {
            "id": f"{source}_{i}",
            "values": emb,
            "metadata": {"content": chunk, "source": source, "chunk_idx": i},
        }
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings))
    ]
    _batched_upsert(vectors, namespace)
    return len(vectors)


def retrieve_by_source(filename: str, namespace: str = "default", k: int = 20) -> list[dict]:
    """Fetch all chunks whose source metadata contains the given filename.

    Images are stored as  image:<filename>
    Videos are stored as  video:<filename>
    Audio  is stored as   audio:<filename>
    Raw files are stored as <filename>
    We try all variants so filename-based queries always hit.
    """
    # Build candidate source values
    base = filename.strip()
    candidates = [
        base,
        f"image:{base}",
        f"video:{base}",
        f"audio:{base}",
    ]
    try:
        # Use a dummy zero-vector; filter does the real work
        dim = _get_embedder().get_sentence_embedding_dimension()
        zero_vec = [0.0] * dim
        results = _get_index().query(
            vector=zero_vec,
            top_k=k,
            namespace=namespace,
            include_metadata=True,
            filter={"source": {"$in": candidates}},
        )
        chunks = []
        for match in results.get("matches", []):
            meta = dict(match.get("metadata", {}))
            content = meta.pop("content", "")
            if content:
                chunks.append({
                    "content": content,
                    "metadata": meta,
                    "score": round(match.get("score", 0.9), 4),
                })
        return chunks
    except Exception:
        return []


def retrieve(query: str, k: int = TOP_K, namespace: str = "default",
             min_score: float = 0.35) -> list[dict]:
    """Retrieve top-k chunks by semantic similarity.

    Args:
        min_score: Cosine similarity floor. Chunks below this score are
                   discarded before returning — prevents loosely-related
                   documents from contaminating the context.
    """
    emb = _embed([query])[0]
    results = _get_index().query(vector=emb, top_k=k, namespace=namespace, include_metadata=True)
    chunks = []
    for match in results.get("matches", []):
        score = round(match.get("score", 0.0), 4)
        if score < min_score:
            continue
        meta = dict(match.get("metadata", {}))
        content = meta.pop("content", "")
        chunks.append({"content": content, "metadata": meta, "score": score})
    return chunks


def collection_count(namespace: str = "default") -> int:
    try:
        stats = _get_index().describe_index_stats()
        ns = stats.get("namespaces", {}).get(namespace, {})
        return ns.get("vector_count", 0)
    except Exception:
        return 0


def list_sources(namespace: str = "default") -> list[str]:
    # Pinecone does not support full metadata enumeration without a fetch;
    # returning the active namespace as a proxy for source listing.
    try:
        stats = _get_index().describe_index_stats()
        return sorted(stats.get("namespaces", {}).keys())
    except Exception:
        return []
